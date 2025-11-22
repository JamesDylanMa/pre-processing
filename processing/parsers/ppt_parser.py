"""
PowerPoint document parser
"""
from typing import Dict, List
from pathlib import Path
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PPTParser:
    """Parse PowerPoint documents"""
    
    def __init__(self):
        self.name = "ppt_parser"
    
    def parse(self, file_path: str) -> Dict:
        """
        Parse PowerPoint file and extract content
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Dictionary with extracted content
        """
        result = {
            "file_path": file_path,
            "parser": self.name,
            "slides": [],
            "text": "",
            "metadata": {}
        }
        
        if not PPTX_AVAILABLE:
            result["error"] = "python-pptx library not available"
            return result
        
        try:
            prs = Presentation(file_path)
            
            result["metadata"]["total_slides"] = len(prs.slides)
            
            all_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_data = {
                    "slide_number": i + 1,
                    "shapes": []
                }
                
                for shape in slide.shapes:
                    shape_data = {
                        "shape_type": shape.shape_type,
                        "text": ""
                    }
                    
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            shape_data["text"] = text
                            slide_text.append(text)
                    
                    slide_data["shapes"].append(shape_data)
                
                slide_data["text"] = "\n".join(slide_text)
                all_text.append(slide_data["text"])
                result["slides"].append(slide_data)
            
            result["text"] = "\n\n".join(all_text)
            
            # Extract core properties if available
            if prs.core_properties:
                result["metadata"]["document_properties"] = {
                    "title": prs.core_properties.title,
                    "author": prs.core_properties.author,
                    "created": str(prs.core_properties.created) if prs.core_properties.created else None,
                    "modified": str(prs.core_properties.modified) if prs.core_properties.modified else None
                }
        
        except Exception as e:
            result["error"] = f"Failed to parse PowerPoint file: {str(e)}"
        
        return result


